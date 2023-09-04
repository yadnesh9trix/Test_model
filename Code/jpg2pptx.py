import os
from pptx import Presentation
# import module
import pandas as pd
from datetime import datetime,timedelta
import warnings
warnings.filterwarnings('ignore')
import os
import warnings
warnings.filterwarnings('ignore')
#-----------------------------------------------------------------------------------------------------------------------
## Define the today's date
today = datetime.today().date()
tday  =today.strftime("%d_%b_%Y")

## Start
if __name__ == '__main__':
    tax_data = "D:/"
    std_path= r"D:\Test_model/"
    inppath = std_path + "Input/"
    outpth = std_path + "OutPut/" + tday + "/"
    paidamount_folder = std_path + "Paidamount/"

    if os.path.exists(outpth):
        pass
    else:
        os.mkdir(outpth)

def convert_images_to_ppt(image_folder, output_file):
    prs = Presentation()
    filenames = ['0.jpg', '1.jpg', '2.jpg', '3.jpg', '4.jpg', '5.jpg', '6.jpg', '7.jpg', '8.jpg', '9.jpg', '10.jpg',
                 '11.jpg', '12.jpg', '13.jpg', '14.jpg', '15.jpg', '16.jpg', '17.jpg', '18.jpg', '19.jpg', '20.jpg',
                 '21.jpg', '22.jpg', '23.jpg', '24.jpg', '25.jpg', '26.jpg', '27.jpg', '28.jpg']
    # for filename in os.listdir(image_folder):
    for filename in filenames:
        if filename.lower().endswith(('.jpg', '.jpeg', '.png', '.gif', '.bmp')):
            img_path = os.path.join(image_folder, filename)
            slide_layout = prs.slide_layouts[6]  # Use slide layout 6 for blank slide
            slide = prs.slides.add_slide(slide_layout)
            left = top = 0
            slide.shapes.add_picture(img_path, left, top, prs.slide_width, prs.slide_height)

    prs.save(output_file)

image_folder = outpth
output_file = outpth + "output_presentation.pptx"
convert_images_to_ppt(image_folder, output_file)